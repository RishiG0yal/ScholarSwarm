"""
PPTX Parser — Extracts text from PowerPoint presentations with slide-level coordinate mapping.
Uses python-pptx to iterate slides and shapes.
"""

import os
import logging
from pptx import Presentation
from backend.models.schemas import TextChunk, ParsedDocument

logger = logging.getLogger("rag_pipeline.parsers.pptx")


def parse_pptx(file_path: str, file_id: str) -> ParsedDocument:
    """
    Parse a PPTX file and return structured text chunks with coordinates.

    Args:
        file_path: Absolute path to the PPTX file.
        file_id: Unique identifier for this file upload.

    Returns:
        ParsedDocument with text chunks mapped to [Slide X, Shape Y].
    """
    filename = os.path.basename(file_path)
    logger.info(f"Parsing PPTX: {filename}")

    prs = Presentation(file_path)
    total_slides = len(prs.slides)
    chunks: list[TextChunk] = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        shape_counter = 0

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            # Collect all text from the shape's text frame
            shape_text_parts = []
            for paragraph in shape.text_frame.paragraphs:
                para_text = paragraph.text.strip()
                if para_text:
                    shape_text_parts.append(para_text)

            if not shape_text_parts:
                continue

            shape_counter += 1
            combined_text = "\n".join(shape_text_parts)

            chunks.append(
                TextChunk(
                    text=combined_text,
                    page_number=slide_num,  # Slide = page
                    paragraph_number=shape_counter,  # Shape = paragraph
                    source_file=filename,
                    chunk_id=f"{file_id}_slide{slide_num}_shape{shape_counter}",
                )
            )

        # Also extract text from tables on the slide
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            table = shape.table
            table_text_parts = []
            for row in table.rows:
                row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_cells:
                    table_text_parts.append(" | ".join(row_cells))

            if table_text_parts:
                shape_counter += 1
                chunks.append(
                    TextChunk(
                        text="\n".join(table_text_parts),
                        page_number=slide_num,
                        paragraph_number=shape_counter,
                        source_file=filename,
                        chunk_id=f"{file_id}_slide{slide_num}_table{shape_counter}",
                    )
                )

    logger.info(f"  Extracted {len(chunks)} chunks from {total_slides} slides")

    return ParsedDocument(
        file_id=file_id,
        filename=filename,
        file_type="pptx",
        total_pages=total_slides,
        total_chunks=len(chunks),
        chunks=chunks,
    )
